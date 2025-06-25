import streamlit as st
import base64
import httpx
from openai import OpenAI
import requests
import os
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE

# --- PAGE CONFIGURATION ---
st.set_page_config(
    page_title="Art Director's Briefing Tool",
    page_icon="🎨",
    layout="wide"
)

# --- SESSION STATE INITIALIZATION ---
if 'slides' not in st.session_state:
    st.session_state.slides = [{
        'id': 0,
        'title': 'Slide 1: Title',
        'text': 'Add your message to users here.',
        'image_prompt': None,
        'image_url': None,
        'image_bytes': None, # To store downloaded image data
        'text_position': 'bottom'
    }]
if 'current_slide_idx' not in st.session_state:
    st.session_state.current_slide_idx = 0
if 'next_id' not in st.session_state:
    st.session_state.next_id = 1
if 'openai_api_key' not in st.session_state:
    st.session_state.openai_api_key = ""

# --- HELPER FUNCTIONS ---

def generate_and_download_image(prompt: str, api_key: str) -> tuple[str, bytes] | None:
    """
    Generates an image using the provided prompt and API key,
    then downloads its content and returns both the URL and the image bytes.
    """
    if not api_key:
        st.error("OpenAI API key is missing. Please enter it in the sidebar.")
        return None
    try:
        # Initialize the OpenAI client correctly as per user's working example
        http_client = httpx.Client()
        client = OpenAI(api_key=api_key, http_client=http_client)
        
        st.info("Generating image with DALL-E 3...")
        response = client.images.generate(
            model="dall-e-3",
            prompt=prompt,
            size="1792x1024", # 16:9 aspect ratio
            quality="standard",
            n=1
        )
        image_url = response.data[0].url
        
        # Download the image content immediately after generation
        st.info("Downloading generated image...")
        image_response = requests.get(image_url, timeout=30)
        image_response.raise_for_status() # Raise an exception for bad status codes
        image_bytes = image_response.content
        st.success("Image successfully generated and downloaded!")
        
        return image_url, image_bytes

    except Exception as e:
        msg = str(e)
        if "Incorrect API key" in msg:
            st.error("The provided OpenAI API key is incorrect. Please check and re-enter it.")
        elif "content_policy_violation" in msg:
            st.error("Image generation failed due to a content policy violation. Please modify your prompt.")
        else:
            st.error(f"Image generation failed: {msg}")
        return None

def create_pptx_from_slides(slides: list[dict]) -> BytesIO:
    """
    Creates a PowerPoint presentation from the slide data, with images as backgrounds
    and formatted text overlays. Returns the presentation as a BytesIO stream.
    """
    prs = Presentation()
    # Use a 16:9 slide layout for widescreen displays
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)

    for slide_data in slides:
        # Add a blank slide layout (layout index 6 is typically blank)
        blank_slide_layout = prs.slide_layouts[6] 
        slide = prs.slides.add_slide(blank_slide_layout)

        # Add image as the full-slide background
        if slide_data.get('image_bytes'):
            image_stream = BytesIO(slide_data['image_bytes'])
            # Add picture to fill the entire slide
            slide.shapes.add_picture(image_stream, Inches(0), Inches(0), width=prs.slide_width, height=prs.slide_height)

        # Add a text box overlay for the user's message
        text = slide_data.get('text', '')
        if text:
            # Define text box dimensions to be slightly smaller than the slide for padding
            tx_width = Inches(14.2)
            tx_height = Inches(2.5) # Generous height to accommodate multiple lines of text
            tx_left = Inches(0.9) # Center the text box horizontally

            # Position the text box based on user's selection ('top', 'center', or 'bottom')
            position = slide_data.get('text_position', 'bottom')
            if position == 'top':
                tx_top = Inches(0.5)
            elif position == 'center':
                tx_top = (prs.slide_height - tx_height) / 2
            else: # Default to 'bottom'
                tx_top = prs.slide_height - tx_height - Inches(0.5)

            # Add the text box shape to the slide
            tx_box = slide.shapes.add_textbox(tx_left, tx_top, tx_width, tx_height)
            
            # --- FORMATTING FOR MAXIMUM READABILITY ---
            # Add a semi-transparent background to the text box
            fill = tx_box.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(0, 0, 0) # Black background
            fill.transparency = 0.30 # 30% transparent so the image is still visible
            
            # Set text frame properties for padding and alignment
            text_frame = tx_box.text_frame
            text_frame.clear() # Clear existing paragraphs
            text_frame.margin_left = Inches(0.2)
            text_frame.margin_right = Inches(0.2)
            text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE # Center text vertically within the box
            text_frame.word_wrap = True
            text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT # Adjust shape size to text
            
            # Add and format the text paragraph
            p = text_frame.paragraphs[0]
            p.text = text
            p.font.name = 'Calibri'
            p.font.size = Pt(32)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255) # White text for high contrast

    # Save the final presentation to a byte stream to be served for download
    pptx_stream = BytesIO()
    prs.save(pptx_stream)
    pptx_stream.seek(0)
    return pptx_stream

# --- UI LAYOUT (3 Columns) ---
col_left, col_center, col_right = st.columns([0.25, 0.45, 0.3])

# --- LEFT SIDEBAR: Settings & Slide Management ---
with col_left:
    st.header("Settings")
    st.session_state.openai_api_key = st.text_input(
        "Enter OpenAI API Key",
        type="password",
        value=st.session_state.openai_api_key,
        help="Your API key is stored temporarily for this session."
    )
    st.write("---")

    st.header("Slides")
    if st.button("➕ Add New Slide", use_container_width=True):
        new_id = st.session_state.next_id
        new_slide = {
            'id': new_id, 'title': f'Slide {new_id + 1}: New Slide', 'text': 'Add your message here.',
            'image_prompt': None, 'image_url': None, 'image_bytes': None, 'text_position': 'bottom'
        }
        st.session_state.slides.append(new_slide)
        st.session_state.next_id += 1
        # Automatically select the new slide
        st.session_state.current_slide_idx = len(st.session_state.slides) - 1
        st.rerun()

    # Display slide list for selection and deletion
    for i, slide in enumerate(list(st.session_state.slides)):
        with st.container(border=True, height=120):
            is_selected = (i == st.session_state.current_slide_idx)
            label = f"Slide {i+1}" + (" (Selected)" if is_selected else "")
            # Button to select a slide for editing
            if st.button(label, key=f"select_{slide['id']}", use_container_width=True, type="primary" if is_selected else "secondary"):
                st.session_state.current_slide_idx = i
                st.rerun()
            # Button to delete a slide
            if st.button("🗑️ Delete", key=f"delete_{slide['id']}", use_container_width=True):
                if len(st.session_state.slides) > 1:
                    st.session_state.slides.pop(i)
                    # Adjust current index if needed to prevent errors
                    if st.session_state.current_slide_idx >= len(st.session_state.slides):
                        st.session_state.current_slide_idx = len(st.session_state.slides) -1
                    st.rerun()
                else:
                    st.warning("Cannot delete the last slide.")

# --- CENTER PANEL: Slide Editor & Preview ---
with col_center:
    st.header("Presentation Editor")
    st.write("---")
    
    current_slide = st.session_state.slides[st.session_state.current_slide_idx]
    
    # Edit fields for the currently selected slide
    current_slide['title'] = st.text_input("Slide Title (for reference)", value=current_slide['title'], key=f"title_{current_slide['id']}")
    current_slide['text'] = st.text_area("Message to Overlay on Image", value=current_slide['text'], height=150, key=f"text_{current_slide['id']}")
    
    pos_options = ['bottom', 'top', 'center']
    try:
        current_pos_index = pos_options.index(current_slide.get('text_position', 'bottom'))
    except ValueError:
        current_pos_index = 0 # Default to bottom if value is invalid
    current_slide['text_position'] = st.selectbox("Text Position on Image", pos_options, index=current_pos_index, key=f"pos_{current_slide['id']}")

    st.write("---")
    st.subheader("Slide Preview")
    
    # Display the image if it exists for the current slide
    if current_slide.get('image_url'):
        st.image(current_slide['image_url'], caption=f"Generated from: {current_slide.get('image_prompt')}")
    else:
        st.info("Use the Image Generator on the right to create a background for this slide.")

# --- RIGHT SIDEBAR: Image Prompt & Generation ---
with col_right:
    st.header("Image Generator")
    st.write("---")
    
    with st.form("prompt_form"):
        st.info("Fill out the variables below to construct the image prompt.")
        subject = st.text_input("Subject", "a silver dragon perched on a jagged cliff")
        action = st.text_input("Action", "roaring toward the stormy sky")
        environment = st.text_input("Environment", "craggy seaside coast at dusk")
        style = st.selectbox("Style", ("digital matte painting, hyper-realistic", "illustration", "abstract", "photorealistic", "cel-shaded anime"))
        # Add more creative fields as needed...
        
        submitted = st.form_submit_button("Generate Image", use_container_width=True, type="primary")
        if submitted:
            if not st.session_state.openai_api_key:
                st.error("Please enter your OpenAI API key in the left sidebar.")
            else:
                final_prompt = f"{subject} {action}, {environment}, in the style of a {style}, 16:9 aspect ratio"
                
                with st.spinner("Generating and downloading your image..."):
                    result = generate_and_download_image(final_prompt, st.session_state.openai_api_key)
                
                if result:
                    image_url, image_bytes = result
                    # Temporarily store the generated data in session state before assigning to a slide
                    st.session_state.generated_image_url = image_url
                    st.session_state.generated_image_bytes = image_bytes
                    st.session_state.generated_prompt = final_prompt
                    st.rerun()

    # This block appears after an image has been successfully generated
    if st.session_state.get("generated_image_url"):
        st.write("---")
        st.subheader("Generated Image")
        st.image(st.session_state.generated_image_url, caption="Generated by DALL-E 3")

        # Provide an immediate download button for the newly generated image
        st.download_button(
            "⬇️ Download This Image", 
            st.session_state.generated_image_bytes, 
            file_name=f"generated_image_{st.session_state.slides[st.session_state.current_slide_idx]['id']}.png", 
            mime="image/png",
            use_container_width=True
        )

        # Button to confirm and use the image for the currently selected slide
        if st.button("✅ Use This Image for Current Slide", use_container_width=True):
            slide = st.session_state.slides[st.session_state.current_slide_idx]
            slide['image_url'] = st.session_state.generated_image_url
            slide['image_bytes'] = st.session_state.generated_image_bytes
            slide['image_prompt'] = st.session_state.generated_prompt
            # Clean up temporary state variables
            del st.session_state.generated_image_url
            del st.session_state.generated_image_bytes
            del st.session_state.generated_prompt
            st.rerun()

    # --- EXPORT SECTION ---
    st.write("---")
    st.header("Export Brief")
    
    # Check if any slides have images before enabling the export button
    can_export = any(s.get('image_bytes') for s in st.session_state.slides)

    # The export button is disabled if no slides have images
    if st.button("Export to PowerPoint (.pptx)", use_container_width=True, disabled=not can_export):
        with st.spinner("Creating your PowerPoint brief..."):
            pptx_bytes = create_pptx_from_slides(st.session_state.slides)
            st.download_button(
                label="⬇️ Download Presentation",
                data=pptx_bytes,
                file_name="Art_Brief.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                use_container_width=True
            )
    elif not can_export:
        st.info("Please add a generated image to at least one slide to enable export.")
